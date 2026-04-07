#!/usr/bin/env python3
"""إنتاج ملف تعريفي احترافي لزيادة سيستم، بستايل قريب من خلفية الموقع، مع تصدير صور جميع الشرائح."""

from __future__ import annotations

import io
import math
import random
from datetime import datetime
from pathlib import Path
from zipfile import ZipFile

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

from PIL import Image, ImageDraw, ImageFilter

SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_DIR = SCRIPT_DIR.parent
OUTPUT_DIR = PROJECT_DIR / "outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

TOKEN_PATH = PROJECT_DIR / "token_docs.json"
SCOPES = ["https://www.googleapis.com/auth/drive.file"]

TITLE = "الملف التعريفي | زيادة سيستم | النسخة الاحترافية"
CONTACT_EMAIL = "ziyadasystem@gmail.com"

ASSETS_DIR = PROJECT_DIR / "assets"
DEPLOY_ZIP = ASSETS_DIR / "ziyada system website deploy-695d3ef3562490d34b06fbbb.zip"
RUNTIME_ASSETS = OUTPUT_DIR / "_deck_assets"
RUNTIME_ASSETS.mkdir(parents=True, exist_ok=True)

PROFILE_IMAGE = RUNTIME_ASSETS / "profile_image.png"
BG_MAIN = RUNTIME_ASSETS / "bg_main_wire.png"
BG_ALT = RUNTIME_ASSETS / "bg_alt_wire.png"
BG_CLOSING = RUNTIME_ASSETS / "bg_closing_wire.png"

BG_DARK = RGBColor(15, 23, 42)
TEXT_MAIN = RGBColor(226, 232, 240)
TEXT_MUTED = RGBColor(180, 196, 220)
BLUE = RGBColor(59, 130, 246)
PURPLE = RGBColor(139, 92, 246)
CYAN = RGBColor(6, 182, 212)
PINK = RGBColor(236, 72, 153)
GLASS = RGBColor(19, 31, 53)

SLIDE_W = 13.333
SLIDE_H = 7.5
SAFE_X = 0.58
SAFE_Y = 0.34


def ensure_assets() -> None:
    if not PROFILE_IMAGE.exists() and DEPLOY_ZIP.exists():
        with ZipFile(DEPLOY_ZIP) as zf:
            if "assets/profile_image.png" in zf.namelist():
                PROFILE_IMAGE.write_bytes(zf.read("assets/profile_image.png"))

    if not BG_MAIN.exists():
        make_wire_background(BG_MAIN, seed=11, tint=(139, 92, 246), star_density=800)
    if not BG_ALT.exists():
        make_wire_background(BG_ALT, seed=29, tint=(59, 130, 246), star_density=900)
    if not BG_CLOSING.exists():
        make_wire_background(BG_CLOSING, seed=73, tint=(6, 182, 212), star_density=850)


def make_wire_background(path: Path, seed: int, tint: tuple[int, int, int], star_density: int) -> None:
    """Generate a text-free background inspired by website Three.js visuals."""
    w, h = 1920, 1080
    rng = random.Random(seed)
    img = Image.new("RGB", (w, h), (10, 25, 62))
    draw = ImageDraw.Draw(img, "RGBA")

    # Tiny stars/particles.
    for _ in range(star_density):
        x = rng.randint(0, w - 1)
        y = rng.randint(0, h - 1)
        s = rng.randint(1, 3)
        a = rng.randint(85, 210)
        draw.rectangle([x, y, x + s, y + s], fill=(59, 130, 246, a))

    # Build knot-like wire loops.
    def loop(cx: float, cy: float, r: float, phase: float, alpha: int) -> list[tuple[float, float]]:
        pts: list[tuple[float, float]] = []
        for i in range(620):
            t = (i / 620.0) * math.pi * 2
            rr = r * (1.0 + 0.17 * math.sin(3.2 * t + phase))
            x = cx + math.cos(t) * rr
            y = cy + math.sin(t) * (rr * 0.82) * (1.0 + 0.08 * math.cos(2.5 * t + phase))
            pts.append((x, y))

        for i in range(len(pts) - 1):
            draw.line([pts[i], pts[i + 1]], fill=(tint[0], tint[1], tint[2], alpha), width=1)
            if i % 14 == 0:
                j = (i + 175) % len(pts)
                draw.line([pts[i], pts[j]], fill=(tint[0], tint[1], tint[2], alpha // 3), width=1)
        return pts

    loop(930, 545, 315, 0.5, 115)
    loop(1090, 510, 270, 2.0, 95)
    loop(770, 520, 255, 3.4, 90)

    # Inner geometric core.
    cx, cy, rad = 960, 540, 115
    points = []
    for i in range(6):
        ang = (math.pi * 2 * i / 6) - 0.1
        points.append((cx + rad * math.cos(ang), cy + rad * math.sin(ang)))
    for i in range(6):
        draw.line([points[i], points[(i + 1) % 6]], fill=(6, 182, 212, 120), width=2)
        draw.line([points[i], (cx, cy)], fill=(59, 130, 246, 70), width=1)

    # Subtle vignette.
    draw.rectangle([0, 0, w, h], fill=(8, 14, 30, 46))
    img = img.filter(ImageFilter.GaussianBlur(0.32))
    img.save(path)


def set_dark_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_bg(slide, bg_path: Path, overlay_alpha: float = 0.35) -> None:
    add_image(slide, bg_path, 0, 0, SLIDE_W, SLIDE_H)
    ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    ov.fill.solid()
    ov.fill.fore_color.rgb = BG_DARK
    ov.fill.transparency = overlay_alpha
    ov.line.fill.background()


def add_progress(slide, step: int, total: int) -> None:
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.16))
    base.fill.solid()
    base.fill.fore_color.rgb = BLUE
    base.line.fill.background()

    width = SLIDE_W * (step / total)
    prog = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(width), Inches(0.16))
    prog.fill.solid()
    prog.fill.fore_color.rgb = PURPLE
    prog.line.fill.background()


def add_wire_badge(slide) -> None:
    ring1 = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(9.15), Inches(1.1), Inches(3.5), Inches(3.5))
    ring1.fill.background()
    ring1.line.color.rgb = BLUE
    ring1.line.width = Pt(2)

    ring2 = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(9.62), Inches(1.57), Inches(2.5), Inches(2.5))
    ring2.fill.background()
    ring2.line.color.rgb = PURPLE
    ring2.line.width = Pt(1.5)

    center = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(10.22), Inches(2.2), Inches(1.13), Inches(1.13))
    center.fill.solid()
    center.fill.fore_color.rgb = CYAN
    center.fill.transparency = 0.35
    center.line.color.rgb = CYAN


def add_header(slide, title: str, subtitle: str, step: int, total: int) -> None:
    add_progress(slide, step, total)

    t = slide.shapes.add_textbox(Inches(SAFE_X), Inches(SAFE_Y + 0.2), Inches(9.0), Inches(0.8))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Noto Kufi Arabic"
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = TEXT_MAIN

    s = slide.shapes.add_textbox(Inches(SAFE_X), Inches(SAFE_Y + 0.9), Inches(9.0), Inches(0.55))
    p2 = s.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Noto Kufi Arabic"
    p2.font.size = Pt(13.6)
    p2.font.color.rgb = TEXT_MUTED


def add_image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    if path.exists():
        try:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
            return
        except Exception:
            pass
    fallback = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fallback.fill.solid()
    fallback.fill.fore_color.rgb = RGBColor(28, 42, 67)
    fallback.fill.transparency = 0.15
    fallback.line.color.rgb = BLUE


def add_card(slide, x: float, y: float, w: float, h: float, title: str, lines: list[str], border: RGBColor = BLUE) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = GLASS
    card.fill.transparency = 0.1
    card.line.color.rgb = border
    card.line.width = Pt(1.2)

    ht = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.14), Inches(w - 0.3), Inches(0.36))
    hp = ht.text_frame.paragraphs[0]
    hp.text = title
    hp.font.name = "Noto Kufi Arabic"
    hp.font.bold = True
    hp.font.size = Pt(14)
    hp.font.color.rgb = TEXT_MAIN

    bt = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.55), Inches(w - 0.33), Inches(h - 0.7))
    tf = bt.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = "Noto Kufi Arabic"
        p.font.size = Pt(11.2)
        p.font.color.rgb = TEXT_MAIN


def add_story_timeline(slide, labels: list[str], y: float = 6.45) -> None:
    x = 0.86
    step = 1.95
    for i, label in enumerate(labels):
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + i * step), Inches(y), Inches(0.23), Inches(0.23))
        node.fill.solid()
        node.fill.fore_color.rgb = CYAN if i % 2 == 0 else PURPLE
        node.line.fill.background()

        txt = slide.shapes.add_textbox(Inches(x + i * step + 0.3), Inches(y - 0.03), Inches(1.7), Inches(0.28))
        p = txt.text_frame.paragraphs[0]
        p.text = label
        p.font.name = "Noto Kufi Arabic"
        p.font.size = Pt(9.4)
        p.font.color.rgb = TEXT_MUTED

        if i < len(labels) - 1:
            ln = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + i * step + 0.23), Inches(y + 0.115),
                Inches(x + (i + 1) * step - 0.04), Inches(y + 0.115),
            )
            ln.line.color.rgb = BLUE


def style_chart(chart) -> None:
    # Avoid Arabic text rendering inside chart bitmap layer.
    chart.has_title = False
    chart.has_legend = False

    if chart.category_axis is not None:
        chart.category_axis.tick_labels.font.size = Pt(7)
        chart.category_axis.tick_labels.font.color.rgb = TEXT_MUTED
    if chart.value_axis is not None:
        chart.value_axis.tick_labels.font.size = Pt(7)
        chart.value_axis.tick_labels.font.color.rgb = TEXT_MUTED
        chart.value_axis.has_major_gridlines = False

    try:
        chart.plot_area.format.fill.solid()
        chart.plot_area.format.fill.fore_color.rgb = RGBColor(14, 25, 44)
    except Exception:
        pass


def add_bar_chart(slide, x: float, y: float, w: float, h: float, cats, before_vals, after_vals) -> None:
    data = CategoryChartData()
    data.categories = cats
    data.add_series("Before", before_vals)
    data.add_series("After", after_vals)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    style_chart(chart)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = PURPLE
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = CYAN


def add_col_chart(slide, x: float, y: float, w: float, h: float, cats, vals) -> None:
    data = CategoryChartData()
    data.categories = cats
    data.add_series("Value", vals)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    style_chart(chart)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE


def add_note(slide, x: float, y: float, w: float, text: str) -> None:
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.32))
    p = t.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Noto Kufi Arabic"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED


def add_metric_row(slide, y: float, values: list[tuple[str, str]]) -> None:
    x = 0.86
    w = 2.86
    h = 1.12
    gap = 0.2
    for i, (title, val) in enumerate(values):
        bx = x + i * (w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = GLASS
        card.fill.transparency = 0.1
        card.line.color.rgb = [BLUE, PURPLE, CYAN, PINK][i % 4]

        t = slide.shapes.add_textbox(Inches(bx + 0.16), Inches(y + 0.1), Inches(w - 0.2), Inches(0.36))
        p = t.text_frame.paragraphs[0]
        p.text = title
        p.font.name = "Noto Kufi Arabic"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED

        v = slide.shapes.add_textbox(Inches(bx + 0.16), Inches(y + 0.46), Inches(w - 0.2), Inches(0.5))
        pv = v.text_frame.paragraphs[0]
        pv.text = val
        pv.font.name = "Inter"
        pv.font.bold = True
        pv.font.size = Pt(18)
        pv.font.color.rgb = TEXT_MAIN


def build_deck(out_path: Path) -> None:
    ensure_assets()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    total = 18

    # 1. الغلاف
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_MAIN, overlay_alpha=0.33)
    add_wire_badge(s)
    add_header(s, "زيادة سيستم", "ملف تعريفي احترافي للشركات المؤسسية في المملكة", 1, total)
    add_card(
        s,
        0.85,
        2.0,
        8.85,
        3.1,
        "خلاصة القيمة",
        [
            "نبني أنظمة تشغيل رقمية تربط الاستراتيجية بالتنفيذ اليومي.",
            "نحوّل العمليات اليدوية إلى تدفقات مؤتمتة قابلة للقياس.",
            "نرفع موثوقية القرار عبر لوحات مؤشرات تشغيلية وتنفيذية.",
        ],
        PURPLE,
    )
    add_story_timeline(s, ["تشخيص", "تصميم", "تنفيذ", "قياس", "تحسين", "توسع"])

    # 2. الملخص التنفيذي
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_ALT, overlay_alpha=0.38)
    add_wire_badge(s)
    add_header(s, "الملخص التنفيذي", "شريك تشغيلي بتركيز مؤسسي ونتائج عملية", 2, total)
    add_metric_row(
        s,
        2.0,
        [
            ("خفض زمن العمليات", "-38%"),
            ("تقليل الأعمال اليدوية", "-46%"),
            ("تحسن انضباط الخدمة", "+31%"),
            ("وضوح القرار التنفيذي", "+42%"),
        ],
    )
    add_card(
        s,
        0.85,
        3.45,
        7.25,
        2.35,
        "التركيز التنفيذي",
        [
            "الحوكمة قبل الأدوات، والقياس قبل التوسع.",
            "تنفيذ تدريجي يحافظ على استمرارية التشغيل.",
            "التعاطف في التصميم لضمان تبني الفرق للحلول.",
        ],
        BLUE,
    )
    add_image(s, PROFILE_IMAGE, 8.35, 3.45, 3.95, 2.35)

    # 3. من نحن
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_MAIN, overlay_alpha=0.4)
    add_wire_badge(s)
    add_header(s, "من نحن", "هوية تشغيلية تجمع التقنية بالإدارة", 3, total)
    add_card(
        s,
        0.85,
        1.9,
        5.9,
        4.45,
        "هوية زيادة سيستم",
        [
            "نخدم الشركات المؤسسية وشركات B2B ذات العمليات المتعددة.",
            "نربط فرق التسويق والمبيعات والتشغيل ضمن نموذج موحد.",
            "نحول البيانات إلى قرارات تنفيذية واضحة.",
            "نلتزم بنتائج قابلة للقياس دون مبالغة دعائية.",
        ],
        BLUE,
    )
    add_card(
        s,
        6.98,
        1.9,
        5.32,
        4.45,
        "قيم التنفيذ",
        [
            "الابتكار: حلول عملية وليست تجريبية.",
            "النزاهة: وضوح كامل في خطة التنفيذ والتكلفة.",
            "التعاطف: تصميم الحل بما يناسب الفريق والعميل.",
            "التميز: جودة تنفيذ قابلة للاستدامة.",
        ],
        PURPLE,
    )

    # 4. التحديات
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_ALT, overlay_alpha=0.42)
    add_wire_badge(s)
    add_header(s, "التحديات التشغيلية", "المشكلات المتكررة داخل المنظومات المؤسسية", 4, total)
    add_card(s, 0.85, 1.9, 3.86, 4.4, "تشتت الأنظمة", ["نقل يدوي بين الأدوات", "ازدواجية بيانات", "تضارب التقارير"], PURPLE)
    add_card(s, 4.95, 1.9, 3.86, 4.4, "بطء القرار", ["انعدام لوحات لحظية", "تأخر قراءة الأداء", "ضعف الاستجابة"], BLUE)
    add_card(s, 9.05, 1.9, 3.25, 4.4, "هدر تشغيلي", ["تكرار المهام", "فرص مبيعات ضائعة", "ارتفاع كلفة الخدمة"], CYAN)

    # 5. المنهج
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_MAIN, overlay_alpha=0.4)
    add_wire_badge(s)
    add_header(s, "منهج التنفيذ", "خارطة عمل من التشخيص إلى التحسين المستمر", 5, total)
    phases = ["تحليل الوضع الحالي", "تصميم النموذج", "تنفيذ وتكامل", "تشغيل تجريبي", "قياس وتحسين", "توسع منضبط"]
    x0 = 0.85
    for i, ph in enumerate(phases):
        che = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x0 + i * 2.04), Inches(2.55), Inches(1.9), Inches(1.5))
        color = [BLUE, PURPLE, CYAN, PINK, BLUE, PURPLE][i]
        che.fill.solid()
        che.fill.fore_color.rgb = color
        che.fill.transparency = 0.22
        che.line.color.rgb = color
        t = s.shapes.add_textbox(Inches(x0 + i * 2.04 + 0.13), Inches(2.98), Inches(1.55), Inches(0.8))
        p = t.text_frame.paragraphs[0]
        p.text = ph
        p.font.name = "Noto Kufi Arabic"
        p.font.bold = True
        p.font.size = Pt(10.6)
        p.font.color.rgb = TEXT_MAIN

    # 6. خدماتنا الأساسية
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_ALT, overlay_alpha=0.4)
    add_wire_badge(s)
    add_header(s, "خدماتنا الأساسية", "الخدمات الرقمية المعتمدة", 6, total)
    add_card(
        s,
        0.85,
        1.95,
        5.82,
        4.5,
        "محور النمو الرقمي",
        [
            "أتمتة الأعمال (Business Automation).",
            "إدارة علاقات العملاء والمبيعات (CRM & Sales).",
            "توليد العملاء B2B.",
            "التسويق الرقمي متعدد القنوات.",
        ],
        BLUE,
    )
    add_card(
        s,
        6.92,
        1.95,
        5.38,
        4.5,
        "محور القنوات والظهور",
        [
            "تطوير المواقع المحولة.",
            "تحسين الظهور في البحث (SEO & SEM).",
            "إدارة وسائل التواصل.",
            "تكامل القنوات مع التتبع والتحليل.",
        ],
        PURPLE,
    )

    # 7. خدماتنا المتقدمة
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_MAIN, overlay_alpha=0.42)
    add_wire_badge(s)
    add_header(s, "خدمات تشغيلية متقدمة", "توسعة الخدمات وفق احتياج التشغيل المؤسسي", 7, total)
    add_card(
        s,
        0.85,
        1.95,
        5.9,
        4.45,
        "التحول التشغيلي",
        [
            "التحول من ملفات Excel إلى نظام متكامل.",
            "لوحات مؤشرات تشغيلية وتنفيذية.",
            "مواءمة مؤشرات التسويق والمبيعات.",
            "حوكمة بيانات العملاء والتجديد والاحتفاظ.",
        ],
        CYAN,
    )
    add_card(
        s,
        6.98,
        1.95,
        5.32,
        4.45,
        "الخدمة والحوكمة",
        [
            "نظام طلبات السكان والصيانة.",
            "تجربة ساكن بمعيار ضيافة.",
            "إدارة الأصول والموردين والعقود.",
            "نظام المعرفة وإجراءات التشغيل القياسية.",
        ],
        PINK,
    )

    # 8. الهندسة التشغيلية
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_ALT, overlay_alpha=0.4)
    add_wire_badge(s)
    add_header(s, "الهندسة التشغيلية للحل", "ترابط الأنظمة والفرق والبيانات في نموذج واحد", 8, total)
    add_card(
        s,
        0.85,
        1.95,
        11.45,
        4.45,
        "طبقات التنفيذ",
        [
            "طبقة القنوات: موقع، إعلانات، نماذج، تواصل.",
            "طبقة إدارة العملاء: CRM مع حوكمة دورة الفرصة.",
            "طبقة الأتمتة: تدفقات عمل تنفذ المهام المتكررة تلقائيا.",
            "طبقة التحليل: لوحات مؤشرات تربط النشاط بالنتائج.",
            "طبقة القرار: إيقاع مراجعة أسبوعي لتنفيذ التحسينات.",
        ],
        BLUE,
    )

    # 9. نموذج الأثر
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_MAIN, overlay_alpha=0.42)
    add_wire_badge(s)
    add_header(s, "نموذج الأثر", "قراءة كمية للفارق قبل وبعد التطبيق", 9, total)
    add_bar_chart(
        s,
        0.85,
        1.95,
        6.15,
        4.15,
        ["Cycle", "Manual", "Decision", "SLA"],
        [100, 100, 100, 100],
        [62, 54, 58, 49],
    )
    add_col_chart(
        s,
        7.3,
        1.95,
        5.0,
        4.15,
        ["Lead-SQL", "SQL-Opp", "Opp-Win"],
        [31, 25, 18],
    )
    add_note(s, 0.95, 6.2, 6.0, "يسار: كفاءة تشغيلية | يمين: تحسن مسار التحويل")

    # 10. مسار العملاء
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_ALT, overlay_alpha=0.42)
    add_wire_badge(s)
    add_header(s, "مسار اكتساب العملاء", "تدفق من الوعي إلى الإغلاق ضمن ضوابط واضحة", 10, total)
    add_col_chart(
        s,
        0.85,
        2.0,
        7.35,
        4.0,
        ["Traffic", "Leads", "MQL", "Opp", "Won"],
        [100, 44, 27, 16, 9],
    )
    add_card(
        s,
        8.45,
        2.0,
        3.85,
        4.0,
        "ضبط الفانل",
        [
            "تأهيل آلي للعملاء المحتملين.",
            "حوكمة تسليم بين التسويق والمبيعات.",
            "متابعة متعددة القنوات.",
            "قياس من الحملة حتى الإيراد.",
        ],
        CYAN,
    )

    # 11. الساكن والصيانة
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_MAIN, overlay_alpha=0.4)
    add_wire_badge(s)
    add_header(s, "خدمات الساكن والصيانة", "انضباط خدمة بمعيار ضيافة", 11, total)
    add_card(
        s,
        0.85,
        1.95,
        6.95,
        4.4,
        "دورة الخدمة",
        [
            "استقبال الطلب وتصنيفه حسب الأولوية.",
            "بدء SLA تلقائيا مع قواعد تصعيد.",
            "توجيه التنفيذ والتحقق من الإغلاق.",
            "جمع تقييم الساكن وقياس الرضا.",
            "تقرير دوري للإدارة والتحسين.",
        ],
        PURPLE,
    )
    add_image(s, PROFILE_IMAGE, 8.05, 1.95, 4.25, 4.4)

    # 12. دراسة حالة 1
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_ALT, overlay_alpha=0.4)
    add_wire_badge(s)
    add_header(s, "دراسة حالة 1", "تحول تشغيلي داخل شركة خدمات", 12, total)
    add_card(s, 0.85, 1.95, 3.95, 4.4, "التحدي", ["تعدد ملفات العمل", "تقارير متضاربة", "تأخر القرار"], BLUE)
    add_card(s, 5.05, 1.95, 3.95, 4.4, "المنهج", ["تصميم نموذج بيانات", "بناء لوحات تنفيذية", "أتمتة دورة التقارير"], CYAN)
    add_card(s, 9.25, 1.95, 3.05, 4.4, "النتيجة", ["خفض زمن التقرير 58%", "رفع الدقة 41%", "قرار أسرع"], PURPLE)

    # 13. دراسة حالة 2
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_MAIN, overlay_alpha=0.42)
    add_wire_badge(s)
    add_header(s, "دراسة حالة 2", "تحسين مسار المبيعات واكتساب العملاء", 13, total)
    add_card(s, 0.85, 1.95, 5.8, 4.4, "قبل التنفيذ", ["ليدز غير مؤهلة", "بطء المتابعة", "غياب قراءة دقيقة للفانل"], BLUE)
    add_card(s, 6.9, 1.95, 5.4, 4.4, "بعد التنفيذ", ["حوكمة MQL/SAL", "متابعة آلية", "تحسن التحويل وقيمة الفرص"], CYAN)

    # 14. الحوكمة والأمان
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_ALT, overlay_alpha=0.4)
    add_wire_badge(s)
    add_header(s, "الأمان والحوكمة", "اعتمادية تشغيلية مناسبة للبيئات المؤسسية", 14, total)
    add_card(s, 0.85, 1.95, 3.85, 4.35, "ضبط الصلاحيات", ["أدوار واضحة", "مسارات اعتماد", "سجلات تدقيق"], BLUE)
    add_card(s, 4.95, 1.95, 3.85, 4.35, "حوكمة البيانات", ["سياسة تعامل", "تتبع تدفق البيانات", "ضبط الاستخدام"], PURPLE)
    add_card(s, 9.05, 1.95, 3.25, 4.35, "استمرارية الأعمال", ["خطط بديلة", "إنذار مبكر", "مراجعة دورية"], CYAN)

    # 15. نماذج التعاقد
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_MAIN, overlay_alpha=0.42)
    add_wire_badge(s)
    add_header(s, "نماذج التعاقد والتنفيذ", "مرونة في البداية والتوسع", 15, total)
    add_card(s, 0.85, 1.95, 3.85, 4.35, "تقييم تشغيلي سريع", ["تشخيص أولي", "تحديد الفجوات", "خارطة أولويات"], BLUE)
    add_card(s, 4.95, 1.95, 3.85, 4.35, "برنامج تحول متكامل", ["تنفيذ شامل", "تكامل الأنظمة", "إدارة التغيير"], PURPLE)
    add_card(s, 9.05, 1.95, 3.25, 4.35, "تشغيل أتمتة مُدار", ["متابعة دورية", "تحسين مستمر", "تقارير تنفيذية"], CYAN)

    # 16. التدريب والتمكين
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_ALT, overlay_alpha=0.42)
    add_wire_badge(s)
    add_header(s, "التدريب والتمكين", "بناء تبني داخلي مستدام", 16, total)
    add_card(
        s,
        0.85,
        1.95,
        11.45,
        4.35,
        "سياسة التدريب",
        [
            "التدريب يتم أونلاين فقط وفق السياسة المعتمدة.",
            "جلسات مباشرة عن بُعد لفرق التشغيل والإدارة.",
            "مكتبة فيديوهات مسجلة للتشغيل اليومي.",
            "متابعة تبني الفريق وربط التدريب بالمؤشرات.",
        ],
        PINK,
    )

    # 17. التواصل والخطوة التالية
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_MAIN, overlay_alpha=0.42)
    add_wire_badge(s)
    add_header(s, "الخطوة التالية", "جلسة تشخيص تنفيذية لتحديد الأولويات", 17, total)
    add_card(
        s,
        0.85,
        2.05,
        10.7,
        3.15,
        "بداية الشراكة",
        [
            "تحديد نطاق المرحلة الأولى خلال 90 يوم.",
            "تثبيت مؤشرات الأداء ومسؤوليات التنفيذ.",
            "البدء بخطة تشغيل واضحة وقابلة للتوسع.",
            f"قناة التواصل الرسمية: {CONTACT_EMAIL}",
        ],
        CYAN,
    )
    add_story_timeline(s, ["اجتماع", "تشخيص", "خطة", "تنفيذ", "قياس"])

    # 18. صفحة الشكر
    s = prs.slides.add_slide(blank)
    set_dark_background(s)
    add_bg(s, BG_CLOSING, overlay_alpha=0.35)
    add_wire_badge(s)
    add_header(s, "شكرا لكم", "نتطلع لبناء أثر تشغيلي مستدام معكم", 18, total)
    add_card(
        s,
        1.1,
        2.2,
        10.95,
        2.8,
        "فريق زيادة سيستم",
        [
            "نشكر وقتكم واهتمامكم.",
            "يسعدنا ترتيب جلسة تنفيذية حسب جدولكم.",
            f"للتواصل: {CONTACT_EMAIL}",
        ],
        PURPLE,
    )

    prs.save(str(out_path))


def get_drive_service():
    if not TOKEN_PATH.exists():
        raise RuntimeError("token_docs.json غير موجود. يلزم توفر توكن Google Drive.")

    creds = Credentials.from_authorized_user_file(str(TOKEN_PATH), SCOPES)
    if creds.expired and creds.refresh_token:
        creds.refresh(Request())
    return build("drive", "v3", credentials=creds)


def upload_as_google_slides(drive, pptx_path: Path) -> dict:
    metadata = {
        "name": f"{TITLE} - {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        "mimeType": "application/vnd.google-apps.presentation",
    }
    media = MediaFileUpload(
        str(pptx_path),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=False,
    )
    return drive.files().create(
        body=metadata,
        media_body=media,
        fields="id,name,mimeType,webViewLink",
    ).execute()


def export_presentation_pdf(drive, file_id: str, out_pdf: Path) -> None:
    req = drive.files().export_media(fileId=file_id, mimeType="application/pdf")
    fh = io.FileIO(out_pdf, "wb")
    downloader = MediaIoBaseDownload(fh, req)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    fh.close()


def export_slide_pngs_from_pdf(pdf_path: Path, out_dir: Path) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        import fitz  # PyMuPDF
    except Exception:
        return 0

    doc = fitz.open(str(pdf_path))
    count = 0
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        target = out_dir / f"slide_{i+1:02d}.png"
        pix.save(str(target))
        count += 1
    doc.close()
    return count


def main():
    out_pptx = OUTPUT_DIR / "ziyada_system_company_profile.pptx"
    build_deck(out_pptx)

    drive = get_drive_service()
    created = upload_as_google_slides(drive, out_pptx)

    link = created.get("webViewLink") or f"https://docs.google.com/presentation/d/{created.get('id')}/edit"
    print("Created file:", created.get("name"))
    print("MIME:", created.get("mimeType"))
    print("Link:", link)

    out_pdf = OUTPUT_DIR / "ziyada_system_company_profile.pdf"
    shots_dir = OUTPUT_DIR / "screenshots"
    try:
        export_presentation_pdf(drive, created.get("id"), out_pdf)
        shots = export_slide_pngs_from_pdf(out_pdf, shots_dir)
        print("PDF:", out_pdf)
        print("Slide screenshots:", shots_dir)
        print("Screenshot count:", shots)
    except Exception as exc:
        print("Screenshot export failed:", exc)


if __name__ == "__main__":
    main()
